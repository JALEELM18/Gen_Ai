import os
import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.llms import OpenAI
import pandas as pd
from pptx import Presentation
from docx import Document
from dotenv import load_dotenv
from time import sleep
from requests.exceptions import HTTPError
import math
import csv
load_dotenv()
class SessionState:
    def __init__(self, **kwargs):
        self.__dict__.update(kwargs)
def process_documents(documents_folder):
    documents = []
    for file_name in os.listdir(documents_folder):
        file_path = os.path.join(documents_folder, file_name)
        if file_path.endswith(".pdf"):
            pdf_reader = PdfReader(file_path)
            for page in pdf_reader.pages:
                documents.append(page.extract_text())
        elif file_path.endswith(".txt"):
            with open(file_path, "r") as file:
                documents.append(file.read())
        elif file_path.endswith(".docx"):
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                documents.append(paragraph.text)
        elif file_path.endswith(".xlsx"):
            xls = pd.ExcelFile(file_path)
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                for col in df.columns:
                    documents.extend(df[col].dropna().astype(str).tolist())
        elif file_path.endswith(".pptx"):
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        documents.append(shape.text)
        elif file_path.endswith(".csv"):  # Handle CSV files
            with open(file_path, 'r', encoding='utf-8') as csvfile:
                reader = csv.reader(csvfile)
                for row in reader:
                    documents.append(' '.join(row))
    text = "\n".join(documents)
    if not text:
        return None
    char_text_splitter = CharacterTextSplitter(separator="\n", chunk_size=1000, chunk_overlap=200, length_function=len)
    text_chunks = char_text_splitter.split_text(text)
    if not text_chunks:
        return None
    embeddings = OpenAIEmbeddings()
    docsearch = FAISS.from_texts(text_chunks, embeddings)
    return docsearch
def answer_query(query, docsearch, chain):
    docs = docsearch.similarity_search(query)
    response = chain.run(input_documents=docs, question=query)
    return response
def delete_files(selected_documents, documents_folder):
    for file_name in selected_documents:
        file_path = os.path.join(documents_folder, file_name)
        try:
            os.remove(file_path)
            st.write(f"{file_name} deleted successfully.")
        except Exception as e:
            st.error(f"Error deleting {file_name}: {e}")
def exponential_backoff_retry(func, max_retries=5, initial_wait_time=1):
    retries = 0
    wait_time = initial_wait_time
    while retries < max_retries:
        try:
            return func()
        except HTTPError as e:
            if e.response.status_code == 429:  # Rate limit exceeded
                st.warning(f"Rate limit exceeded. Retrying in {wait_time} seconds...")
                sleep(wait_time)
                wait_time *= 2  # Exponential backoff
                retries += 1
            else:
                raise  # Reraise the exception if it's not a rate limit error
    st.error("Max retry attempts reached. Please try again later.")
def main():
    st.header("Document Search \U0001F916")  # Unicode for robot emoji
    session_state = SessionState(documents_folder=None, docsearch=None, response="", is_admin=False)
    chain = None  # Initialize chain variable outside of the condition
    # Authentication mechanism
    user_type = st.selectbox("Select user type:", ["Admin", "User"])
    if user_type == "Admin":
        # Admin email and password authentication
        admin_email = st.text_input("Enter admin email:")
        admin_password = st.text_input("Enter admin password:", type="password")
        if admin_email == "jaleel@gmail.com" and admin_password == "jaleel@786":
            session_state.is_admin = True
        else:
            st.error("Incorrect email or password. Please try again.")
            return
        if session_state.is_admin:
            # Admin functionality
            uploaded_files = st.file_uploader("Upload your documents", accept_multiple_files=True, type=["pdf", "txt", "docx", "xlsx", "pptx", "csv"])
            documents_folder = "uploaded_documents_admin"
            if not os.path.exists(documents_folder):
                os.makedirs(documents_folder)
            for file in uploaded_files:
                with open(os.path.join(documents_folder, file.name), "wb") as f:
                    f.write(file.getvalue())
            session_state.documents_folder = documents_folder
            session_state.docsearch = process_documents(documents_folder)
            openai_api_key = st.text_input("Enter your OpenAI API key:", type="password")  # Ask for API key
            def load_chain():
                return load_qa_chain(OpenAI(api_key=openai_api_key), chain_type="stuff")
            chain = exponential_backoff_retry(load_chain)
            # Display uploaded documents
            st.subheader("Uploaded Documents")
            selected_documents = []
            for file_name in os.listdir(documents_folder):
                selected = st.checkbox(file_name)
                if selected:
                    selected_documents.append(file_name)
            # Delete selected documents
            if st.button("Delete Selected Documents"):
                delete_files(selected_documents, documents_folder)
               
    else:
        # Normal user functionality
        documents_folder = "uploaded_documents_admin"  # Use the admin's uploaded documents
        if not os.path.exists(documents_folder):
            st.write("No documents uploaded yet. Please wait for the admin to upload documents.")
            return
 
 
        session_state.documents_folder = documents_folder
        session_state.docsearch = process_documents(documents_folder)
 
        # Display uploaded documents for user
        st.subheader("Uploaded Documents")
        for file_name in os.listdir(documents_folder):
            st.write(file_name)
 
        openai_api_key = os.getenv("OPENAI_API_KEY")  # Get OpenAI API key from environment variable
 
        def load_chain():
            return load_qa_chain(OpenAI(api_key=openai_api_key), chain_type="stuff")
 
        chain = exponential_backoff_retry(load_chain)
 
    # AI assistant chat interface
    st.sidebar.header("Document Query Box 💬")
    chat_history = st.sidebar.empty()  # Placeholder for chat history
    user_input = st.sidebar.text_input("Query:", key="user_input", placeholder="Enter your question here")
    if st.sidebar.button("Submit"):
        if chain is not None:  # Ensure chain is initialized before using it
            response = answer_query(user_input, session_state.docsearch, chain)
            st.sidebar.text("Result:")
            st.sidebar.text(response)

if __name__ == "__main__":
    main()